first_slide_template ="""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Create a new presentation object
prs = Presentation()

# Add a blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Set the background color of the slide to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(128, 128, 128)

# Add image to the left side of the slide with a margin at the top and bottom
left = Inches(0)
top = Inches(0)
height = prs.slide_height
width = prs.slide_width * 3/5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# Add title text box positioned higher
left = prs.slide_width * 3/5
top = Inches(0)
width = prs.slide_width * 2/5
height = prs.slide_height / 2
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(38)
title_p.font.color.rgb = RGBColor(0, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Add subtitle text box
left = prs.slide_width * 3/5
top = prs.slide_height / 2
width = prs.slide_width * 2/5
height = prs.slide_height / 2
subtitle_box = slide.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = subtitle_text
subtitle_p.font.size = Pt(22)
subtitle_p.font.color.rgb = RGBColor(0, 100, 200)
subtitle_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
"""

text_template="""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Create a new presentation object
prs = Presentation()

# Add a blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Set the background color of the slide to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(128, 128, 128)


# Add title text box positioned higher
left = 0
top = Inches(0)
width = prs.slide_width
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = slide_title
title_p.font.bold = True
title_p.font.size = Pt(22)
title_p.font.color.rgb = RGBColor(0, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# Add hardcoded "Remarks" text and bullet points
left = Inches(0.2)
top = Inches(1.5)
width = prs.slide_width - 2*left
height = prs.slide_height * 10/16
insights_box = slide.shapes.add_textbox(left, top, width, height)
insights_frame = insights_box.text_frame
insights_p = insights_frame.add_paragraph()
insights_p.text = slide_theme
insights_p.font.bold = True
insights_p.font.size = Pt(24)
insights_p.font.color.rgb = RGBColor(0, 200, 128)
insights_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


text_p = insights_frame.add_paragraph()
text_p.text = slide_text
text_p.font.size = Pt(14)
text_p.font.color.rgb = RGBColor(255, 255, 255)
text_p.line_spacing = 1.5
text_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
"""

context_template="""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


# Create a new presentation object
prs = Presentation()

# Add a blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Set the background color of the slide to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(128, 128, 128)

# Add image placeholder on the left side of the slide
left = Inches(0.2)
top = Inches(1.8)
height = prs.slide_height - Inches(3)
width = prs.slide_width * 3/5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# Add title text spanning the whole width
left = Inches(0)
top = Inches(0)
width = prs.slide_width
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.margin_top = Inches(0.1)
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(28)
title_p.font.color.rgb = RGBColor(0, 0, 200)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Add hardcoded "Remarks" text and bullet points
left = prs.slide_width * 2/3
top = Inches(1.5)
width = prs.slide_width * 1/3
height = Inches(4.5)
insights_box = slide.shapes.add_textbox(left, top, width, height)
insights_frame = insights_box.text_frame
insights_p = insights_frame.add_paragraph()
insights_p.text = "Remarks:"
insights_p.font.bold = True
insights_p.font.size = Pt(24)
insights_p.font.color.rgb = RGBColor(0, 0, 100)
insights_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
insights_frame.add_paragraph()


# Add bullet points with individual formatting
bullet_points_list = bullet_points.split('\n')
for point in bullet_points_list:
    if point.strip():
        bullet_p = insights_frame.add_paragraph()
        bullet_p.text = point.strip()
        bullet_p.font.size = Pt(12)
        bullet_p.font.color.rgb = RGBColor(0, 0, 200)
        bullet_p.line_spacing = 1.5
        bullet_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
"""


thanks_template=""""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Create a new presentation object
prs = Presentation()

title_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Thanks"
title.font_size = Pt(40)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Set the background color of the slide to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(128, 128, 128)

# Center the title horizontally and vertically
left = int(prs.slide_width / 2.5)
# width = prs.slide_width  // 2
# height = prs.slide_height // 2
top = (prs.slide_height - title.height) // 2
title.left = left
title.top = top
title.width = Inches(2)
title.height = Inches(1)
title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
"""